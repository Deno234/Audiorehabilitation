from __future__ import annotations

import shutil
import struct
import wave
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt


ROOT = Path("/mnt/d/audiorehabilitation")
SOURCE_PPTX = Path("/mnt/c/Users/Gaming 3/Desktop/ROGJ_Tema16_Audiorehabilitacija_prezentacija.pptx")
OUT_DIR = ROOT / "outputs/manual-20260613-ppt/presentations/audiorehab-results/output"
PPTX_OUT = OUT_DIR / "ROGJ_Tema16_Audiorehabilitacija_rezultati_v8.pptx"
NOTES_OUT = OUT_DIR / "ROGJ_Tema16_Audiorehabilitacija_govorne_biljeske_v8.md"
AUDIO_DEMO_DIR = OUT_DIR / "audio_demo"


BLUE = RGBColor(31, 93, 122)
TEAL = RGBColor(34, 161, 178)
DARK = RGBColor(42, 52, 61)
MUTED = RGBColor(96, 108, 117)
LIGHT = RGBColor(241, 246, 248)
PALE = RGBColor(224, 239, 244)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(79, 148, 93)
ORANGE = RGBColor(216, 141, 61)
RED = RGBColor(183, 79, 79)
PURPLE = RGBColor(105, 91, 166)


def delete_all_slides(prs: Presentation) -> None:
    xml_slides = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public delete API.
    for sld_id in list(xml_slides):
        r_id = sld_id.rId
        prs.part.drop_rel(r_id)
        xml_slides.remove(sld_id)


def add_textbox(
    slide,
    text: str,
    x,
    y,
    w,
    h,
    *,
    font_size=18,
    color=DARK,
    bold=False,
    align=PP_ALIGN.LEFT,
    font="Aptos",
    valign=MSO_ANCHOR.TOP,
    margin=0.05,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Cm(margin)
    tf.margin_right = Cm(margin)
    tf.margin_top = Cm(margin)
    tf.margin_bottom = Cm(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_multiline(
    slide,
    lines: list[str],
    x,
    y,
    w,
    h,
    *,
    font_size=16,
    color=DARK,
    first_bold=False,
    bullet=False,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Cm(0.08)
    tf.margin_right = Cm(0.08)
    tf.margin_top = Cm(0.04)
    tf.margin_bottom = Cm(0.04)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(3)
        p.line_spacing = line_spacing
        if bullet and idx > 0:
            p.text = f"• {line}"
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(font_size if idx else font_size + (1 if first_bold else 0))
            run.font.bold = bool(first_bold and idx == 0)
            run.font.color.rgb = color
    return box


def add_title(slide, title: str, slide_no: int, subtitle: str = "Računalna obrada govora i jezika · diplomski studij") -> None:
    add_textbox(slide, title, Cm(0.0), Cm(0.68), Cm(27.2), Cm(1.0), font_size=20, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, f"{slide_no:02d}", Cm(24.65), Cm(0.46), Cm(1.2), Cm(0.55), font_size=13, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, subtitle, Cm(1.2), Cm(18.05), Cm(14.0), Cm(0.5), font_size=8.5, color=MUTED)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0.0), Cm(1.92), Cm(27.2), Cm(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = PALE
    line.line.fill.background()


def add_card(slide, x, y, w, h, title: str, body: str, *, accent=TEAL, title_size=15, body_size=13):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = PALE
    shape.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Cm(0.18), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_textbox(slide, title, x + Cm(0.35), y + Cm(0.25), w - Cm(0.65), Cm(0.6), font_size=title_size, color=BLUE, bold=True)
    add_multiline(slide, body.split("\n"), x + Cm(0.35), y + Cm(0.95), w - Cm(0.65), h - Cm(1.05), font_size=body_size, color=DARK)


def add_metric(slide, x, y, w, h, value: str, label: str, *, color=TEAL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = PALE
    add_textbox(slide, value, x, y + Cm(0.2), w, Cm(0.65), font_size=19, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, x + Cm(0.1), y + Cm(0.95), w - Cm(0.2), h - Cm(0.95), font_size=9.8, color=DARK, align=PP_ALIGN.CENTER)


def add_table(slide, x, y, w, h, headers: list[str], rows: list[list[str]], widths: list[float] | None = None, font_size=8.4):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h)
    table = table_shape.table
    if widths:
        for i, rel in enumerate(widths):
            table.columns[i].width = int(w * rel)
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = "Aptos"
                run.font.size = Pt(font_size)
                run.font.bold = True
                run.font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = DARK
    return table_shape


def add_bar_chart(slide, x, y, w, h, categories: list[str], series: list[tuple[str, list[float]]], *, title: str = ""):
    data = CategoryChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, data)
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 100
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(8)
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
    colors = [TEAL, ORANGE, BLUE, PURPLE]
    for idx, ser in enumerate(chart.series):
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = colors[idx % len(colors)]
    return chart_shape


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_textbox(slide, "Evaluacija jezičnih generativnih modela\nza potrebe audiorehabilitacije", Cm(0.0), Cm(2.7), Cm(27.2), Cm(2.1), font_size=28, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    labels = ["LLM", "fonemi", "validacija", "TTS", "audio"]
    for i, lab in enumerate(labels):
        x = Cm(2.0 + i * 4.9)
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, Cm(8.4), Cm(2.0), Cm(2.0))
        circ.fill.solid()
        circ.fill.fore_color.rgb = PALE if i % 2 else LIGHT
        circ.line.color.rgb = TEAL
        add_textbox(slide, lab, x - Cm(0.4), Cm(10.55), Cm(2.8), Cm(0.55), font_size=13, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Denis Ibiši", Cm(22.0), Cm(17.35), Cm(4.5), Cm(0.55), font_size=12, color=MUTED, align=PP_ALIGN.RIGHT)
    return slide


def amplify_wav(src: Path, dst: Path, target_peak: int = 26000) -> None:
    """Create a louder demo WAV without changing the research source files."""
    with wave.open(str(src), "rb") as wav:
        params = wav.getparams()
        frames = wav.readframes(wav.getnframes())
    if params.sampwidth != 2:
        shutil.copyfile(src, dst)
        return
    sample_count = len(frames) // 2
    samples = struct.unpack(f"<{sample_count}h", frames)
    peak = max((abs(sample) for sample in samples), default=0)
    if peak <= 0:
        shutil.copyfile(src, dst)
        return
    factor = min(target_peak / peak, 8.0)
    louder_samples = [
        max(-32768, min(32767, int(sample * factor)))
        for sample in samples
    ]
    louder = struct.pack(f"<{sample_count}h", *louder_samples)
    with wave.open(str(dst), "wb") as wav:
        wav.setparams(params)
        wav.writeframes(louder)


def prepare_audio_demo() -> dict[str, Path]:
    AUDIO_DEMO_DIR.mkdir(parents=True, exist_ok=True)
    base = ROOT / "outputs/audio_comparison/task16_chatgpt_hjp_tts_subset"
    candidate_id = "20260604_160319_00370"
    sources = {
        "eSpeak NG": base / f"espeak_ng/{candidate_id}.wav",
        "Coqui VITS HR": base / f"coqui_vits_hr/{candidate_id}.wav",
        "SpeechT5 HR": base / f"speecht5_hr/{candidate_id}.wav",
    }
    outputs: dict[str, Path] = {}
    for label, src in sources.items():
        slug = label.lower().replace(" ", "_")
        dst = AUDIO_DEMO_DIR / f"demo_draga_rada_radi_{slug}.wav"
        if src.exists():
            amplify_wav(src, dst)
        outputs[label] = dst
    readme = AUDIO_DEMO_DIR / "README.txt"
    readme.write_text(
        "Audio demo za prezentaciju.\n"
        "Sva tri WAV-a koriste isti tekst: Draga rada radi.\n"
        "Datoteke su pojačane samo za prezentacijsko slušanje; originalni istraživački WAV-ovi nisu mijenjani.\n",
        encoding="utf-8",
    )
    return outputs


def build_deck() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    audio_demo = prepare_audio_demo()
    prs = Presentation(str(SOURCE_PPTX))
    delete_all_slides(prs)
    prs.slide_width = Cm(27.2)
    prs.slide_height = Cm(15.3)

    title_slide(prs)

    # 02 Motivation
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Motivacija: zašto baš fonetski kontroliran materijal?", 2)
    add_card(s, Cm(1.3), Cm(4.2), Cm(7.1), Cm(4.0), "Slušni problem", "U audiorehabilitaciji često se vježbaju kontrasti koje korisnik teško razlikuje.\nPrimjer: /s/ i /š/, /č/ i /ć/, /nj/ i /lj/.", accent=TEAL)
    add_card(s, Cm(10.0), Cm(4.2), Cm(7.1), Cm(4.0), "Kontrolirani tekst", "Treba znati koliko je ciljanih fonema u riječi ili rečenici.\nZato uvodimo klase fonema i postotak zasićenja.", accent=ORANGE)
    add_card(s, Cm(18.7), Cm(4.2), Cm(7.1), Cm(4.0), "Audio vježbe", "Tekst se smije pretvoriti u zvuk tek nakon validacije.\nTTS je zadnji korak, ne dokaz kvalitete.", accent=GREEN)

    # 03 Research question
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Istraživačka pitanja i doprinosi", 3)
    add_metric(s, Cm(1.2), Cm(3.0), Cm(5.1), Cm(2.35), "1", "LLM generiranje: mogu li modeli predložiti hrvatske riječi i rečenice s ciljanim fonemskim zasićenjem?", color=TEAL)
    add_metric(s, Cm(7.4), Cm(3.0), Cm(5.1), Cm(2.35), "2", "Python validacija: mogu li se fonemi, pragovi, duplikati i greške provjeriti deterministički?", color=ORANGE)
    add_metric(s, Cm(13.6), Cm(3.0), Cm(5.1), Cm(2.35), "3", "TTS evaluacija: mogu li validirani kandidati postati tehnički ispravni hrvatski WAV zapisi?", color=GREEN)
    add_metric(s, Cm(19.8), Cm(3.0), Cm(5.1), Cm(2.35), "4", "ASR + slušanje: koliko su sintetizirani zapisi razumljivi i prirodni?", color=PURPLE)
    add_card(s, Cm(2.0), Cm(7.2), Cm(23.0), Cm(4.0), "Glavni doprinos", "Reproducibilni eksperimentalni pipeline: LLM/CSV kandidati → hrvatski fonemizator → deterministička validacija → Hunspell/HJP review → TTS usporedba → ASR WER/CER + slušna provjera.", accent=BLUE, title_size=16, body_size=13.5)

    # 04 Paper relation
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Veza s radom Andrijašević i Vukelić (MIPRO 2024)", 4)
    add_card(s, Cm(1.0), Cm(3.0), Cm(7.5), Cm(4.5), "Što preuzimamo", "• hrvatski materijal za auditivni trening\n• pet fonemskih klasa\n• zasićenje ciljnom klasom\n• riječi i kratke rečenice", accent=TEAL, body_size=12.5)
    add_card(s, Cm(9.7), Cm(3.0), Cm(7.5), Cm(4.5), "Što mijenjamo", "• validacija nije u rukama LLM-a\n• sve ide kroz Python\n• spremamo sirove odgovore, CSV i izvještaje\n• dodajemo PCD, TTS, ASR i slušnu provjeru", accent=ORANGE, body_size=12.5)
    add_card(s, Cm(18.4), Cm(3.0), Cm(7.5), Cm(4.5), "Što uspoređujemo", "• reprodukcija paper-style promptova\n• ChatGPT Plus vs. Ollama\n• paper_style vs. strict_plain_list\n• eSpeak, Coqui i SpeechT5 za audio", accent=GREEN, body_size=12.5)
    add_textbox(s, "Važno: HJP u promptu nije dokaz. HJP valjanost je odvojena ručna / riječ-po-riječ provjera.", Cm(2.0), Cm(9.6), Cm(23.0), Cm(1.0), font_size=14, color=RED, bold=True, align=PP_ALIGN.CENTER)

    # 05 phoneme classes
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Fonemski zahtjevi: pet klasa i formula zasićenja", 5)
    rows = [
        ["N / Niski", "m, n, nj, b, p, u"],
        ["SN / Srednjeniski", "v, g, o, h, l, lj"],
        ["S / Srednji", "a, k, r, d, dž, f, ž"],
        ["SV / Srednjevisoki", "č, e, š, t, đ, j"],
        ["V / Visoki", "ć, i, c, z, s"],
    ]
    add_table(s, Cm(1.3), Cm(3.0), Cm(14.6), Cm(5.0), ["Klasa", "Fonemi"], rows, widths=[0.32, 0.68], font_size=11)
    add_card(s, Cm(17.0), Cm(3.0), Cm(8.7), Cm(5.0), "Formula", "zasićenje = broj fonema ciljne klase / ukupan broj fonema × 100\n\nParser najprije prepoznaje dž, lj i nj kao jedan fonem.", accent=BLUE, body_size=13.5)
    add_card(s, Cm(4.0), Cm(9.3), Cm(19.4), Cm(2.0), "Primjer", "„puno banana u panju” → p u n o b a n a n a u p a nj u\n`nj` se broji kao jedan fonem, razmaci se ne broje.", accent=TEAL, body_size=13)

    # 06 pipeline
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Implementirani postupak: generiranje + deterministička validacija", 6)
    steps = [
        ("1", "prompt", "klasa, zasićenje, tip teksta"),
        ("2", "LLM", "ChatGPT Plus ili Ollama generira kandidate"),
        ("3", "fonemizator", "normalizacija + dž/lj/nj longest-match"),
        ("4", "validacija", "SL, znakovi, 3–5 riječi, duplikati"),
        ("5", "leksika", "Hunspell + ručni HJP word-review"),
        ("6", "izvještaj", "CSV, Markdown, tablice, greške, PCD"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        x = Cm(0.8 + i * 4.25)
        add_metric(s, x, Cm(4.0), Cm(3.55), Cm(2.4), num, title, color=TEAL if i < 3 else ORANGE)
        add_textbox(s, desc, x, Cm(6.65), Cm(3.55), Cm(1.1), font_size=9.5, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(s, "Sva čitanja i pisanja CSV-a rade u UTF-8, tako da se čuvaju č, ć, đ, š, ž i fonemi dž/lj/nj.", Cm(2.0), Cm(10.0), Cm(23.0), Cm(0.8), font_size=13, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    # 07 technical implementation
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Tehnička implementacija: moduli i podaci", 7)
    rows = [
        ["phonemizer.py", "normalizacija, UTF-8, dž/lj/nj longest-match", "tekst → fonemi"],
        ["phoneme_classes.py", "mapiranje fonema u N, SN, S, SV, V", "fonem → klasa"],
        ["validators.py", "saturation, znakovi, broj riječi, duplikati, Hunspell", "candidate → pass/fail + reasons"],
        ["metrics.py", "PCD, frekvencije fonema, duplicate rate", "validirani redovi → metrike"],
        ["generators.py", "manual CSV, Ollama parser, prompt strategije", "LLM output → CSV kandidati"],
        ["report.py / pipeline.py", "run_id, config snapshot, CSV/Markdown report", "reproducibilni eksperimenti"],
        ["tts.py / asr_eval.py", "WAV normalizacija, TTS manifest, WER/CER", "validirani tekst → audio evaluacija"],
    ]
    add_table(s, Cm(1.0), Cm(2.6), Cm(25.2), Cm(7.3), ["Modul", "Što radi", "Ulaz/izlaz"], rows, widths=[0.24, 0.46, 0.30], font_size=7.8)
    add_card(s, Cm(1.5), Cm(10.7), Cm(24.2), Cm(1.6), "Reproducibilnost", "Svaki run ima `run_id`, kopiju konfiguracije, UTF-8 CSV izlaze, sirove generacije i Markdown izvještaj. Zato se rezultati mogu ponovno provjeriti bez ručnog brojanja.", accent=BLUE, body_size=11.5)

    # 08 algorithm
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Ključni algoritam: fonemizacija i validacija", 8)
    add_card(s, Cm(1.0), Cm(2.7), Cm(7.7), Cm(6.3), "Fonemizacija", "1. lowercase\n2. ukloni interpunkciju\n3. očuvaj č, ć, đ, š, ž\n4. longest-match: dž, lj, nj\n5. ostala slova kao pojedinačni fonemi\n6. razmaci se ne broje", accent=TEAL, body_size=11.3)
    add_card(s, Cm(9.7), Cm(2.7), Cm(7.7), Cm(6.3), "Saturation", "target_count = broj fonema iz ciljne klase\nN = ukupan broj fonema\n\nsaturation = target_count / N × 100\n\npasses = saturation ≥ threshold", accent=ORANGE, body_size=11.3)
    add_card(s, Cm(18.4), Cm(2.7), Cm(7.7), Cm(6.3), "Failure reasons", "Svaki kandidat može imati više razloga pada:\n• failed_saturation\n• invalid_characters\n• wrong_word_count\n• duplicate\n• repeated_words\n• dictionary_failed", accent=RED, body_size=11.3)
    add_textbox(s, "Primjer: `panj` → p, a, nj; `džep` → dž, e, p. Ovo je deterministički parser, ne LLM procjena.", Cm(1.8), Cm(10.4), Cm(23.6), Cm(0.9), font_size=13, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    # 09 design
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Dizajn eksperimenata", 9)
    add_card(s, Cm(1.2), Cm(3.0), Cm(7.4), Cm(5.2), "Reprodukcija rada", "Riječi: 5 klasa × 5 SL × 11 = 275\nRečenice: 5 klasa × 2 SL × 21 = 210\nPrompt: paper_style\nIzvor: ChatGPT Plus", accent=TEAL, body_size=13)
    add_card(s, Cm(9.9), Cm(3.0), Cm(7.4), Cm(5.2), "Task 16 usporedba", "5 klasa × 2 SL × 2 tipa × 2 prompta × 20\nChatGPT Plus: 797 kandidata\nOllama: 757 kandidata", accent=ORANGE, body_size=13)
    add_card(s, Cm(18.6), Cm(3.0), Cm(7.4), Cm(5.2), "Audio faza", "TTS-ready skup: 355 kandidata\nUravnoteženi TTS podskup: 96\nSinteza: eSpeak NG, Coqui VITS HR, SpeechT5 HR", accent=GREEN, body_size=13)
    add_textbox(s, "SL = saturation level / prag zasićenja ciljnom fonemskom klasom.", Cm(3.0), Cm(9.7), Cm(21.0), Cm(0.8), font_size=12, color=MUTED, align=PP_ALIGN.CENTER)

    # 10 paper reproduction, words
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Reprodukcija rada: tablica za riječi kao u paperu", 10)
    paper_words = [
        ["L", "100", "63,64", "27,27", "18,18", "54,55"],
        ["LM", "90,91", "81,82", "63,64", "9,09", "0"],
        ["M", "100", "72,73", "72,73", "54,55", "63,64"],
        ["MH", "100", "81,82", "45,45", "36,36", "18,18"],
        ["H1", "36,36", "27,27", "36,36", "0", "0"],
        ["H2", "18,18", "18,18", "0", "0", "9,09"],
    ]
    our_words = [
        ["N", "81,82", "100", "100", "90,91", "90,91"],
        ["SN", "90,91", "90,91", "90,91", "90,91", "90,91"],
        ["S", "100", "100", "100", "90,91", "100"],
        ["SV", "100", "100", "81,82", "36,36", "72,73"],
        ["V", "100", "100", "81,82", "100", "36,36"],
    ]
    add_table(s, Cm(1.0), Cm(3.0), Cm(12.1), Cm(5.2), ["Paper\nklasa", "40 %", "50 %", "60 %", "70 %", "80 %"], paper_words, widths=[0.22, 0.156, 0.156, 0.156, 0.156, 0.156], font_size=8.2)
    add_table(s, Cm(14.1), Cm(3.0), Cm(12.1), Cm(4.5), ["Naša\nklasa", "40 %", "50 %", "60 %", "70 %", "80 %"], our_words, widths=[0.22, 0.156, 0.156, 0.156, 0.156, 0.156], font_size=8.2)
    add_card(s, Cm(1.5), Cm(9.2), Cm(24.2), Cm(2.7), "Isti indikator za riječi", "Paper Table II prikazuje postotak riječi koje zadovoljavaju oba kriterija: saturation level + HJP/standardni jezik. Naša analogna tablica koristi saturation pass + ručni HJP word-review. Duplikati su dodatno analizirani u našem pipelineu, ali nisu dio ove paper-style ćelije.", accent=BLUE, body_size=11.8)

    # 11 paper reproduction, sentences
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Reprodukcija rada: tablica za rečenice kao u paperu", 11)
    paper_sent = [
        ["L", "95,24", "0"],
        ["LM", "100", "95,24"],
        ["M", "95,24", "33,33"],
        ["MH", "95,24", "0"],
        ["H", "66,67", "100"],
    ]
    our_sent = [
        ["N", "100", "76,19"],
        ["SN", "100", "90,48"],
        ["S", "100", "95,24"],
        ["SV", "100", "85,71"],
        ["V", "85,71", "33,33"],
    ]
    add_table(s, Cm(2.0), Cm(3.1), Cm(9.8), Cm(4.6), ["Paper\nklasa", "50 %", "70 %"], paper_sent, widths=[0.42, 0.29, 0.29], font_size=10)
    add_table(s, Cm(15.4), Cm(3.1), Cm(9.8), Cm(4.6), ["Naša\nklasa", "50 %", "70 %"], our_sent, widths=[0.42, 0.29, 0.29], font_size=10)
    add_card(s, Cm(1.6), Cm(8.8), Cm(24.0), Cm(2.9), "Isti indikator za rečenice", "Paper Table III prikazuje postotak rečenica koje zadovoljavaju saturation criterion. Naša analogna tablica koristi isti saturation-pass kriterij. Rezultat je vrlo jak na 50 %, a na 70 % jasno pada za klasu V.", accent=ORANGE, body_size=12)

    # 12 ChatGPT vs Ollama
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Rezultat 2: ChatGPT Plus vs. lokalni Ollama", 12)
    add_bar_chart(s, Cm(1.4), Cm(3.0), Cm(12.5), Cm(6.2), ["ChatGPT Plus", "Ollama"], [("tehnički valjano", [48.6, 4.5]), ("SL prolaz", [85.9, 4.6]), ("Hunspell", [95.4, 76.8])], title="Ukupna usporedba (%)")
    rows = [
        ["ChatGPT Plus", "797", "48,6 %", "duplikati=361; SL=112"],
        ["Ollama llama3.1:8b", "757", "4,5 %", "SL=722; broj riječi=195"],
    ]
    add_table(s, Cm(15.0), Cm(3.35), Cm(10.8), Cm(3.2), ["Izvor", "N", "valid", "glavne greške"], rows, widths=[0.34, 0.14, 0.17, 0.35], font_size=8.5)
    add_card(s, Cm(15.0), Cm(7.3), Cm(10.8), Cm(2.8), "Interpretacija", "ChatGPT Plus najčešće razumije fonemski zadatak, ali se ponavlja. Ollama uglavnom ne pogađa ciljno zasićenje, pa nije dovoljno pouzdan u ovoj konfiguraciji.", accent=ORANGE, body_size=12.3)

    # 13 Prompt strategy/failures
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Što je najčešće pošlo krivo?", 13)
    add_card(s, Cm(1.2), Cm(3.0), Cm(7.6), Cm(5.0), "ChatGPT Plus", "Najveći problem: duplikati.\n\n• strict_plain_list za riječi: 188 duplikata\n• paper_style za riječi: 71,5 % tehnički valjano\n• saturacija uglavnom dobra", accent=TEAL, body_size=12.5)
    add_card(s, Cm(9.8), Cm(3.0), Cm(7.6), Cm(5.0), "Ollama", "Najveći problem: fonetski kriterij.\n\n• 722 kandidata pala su na zasićenju\n• 195 pogrešan broj riječi\n• lokalni model ne broji foneme pouzdano", accent=RED, body_size=12.5)
    add_card(s, Cm(18.4), Cm(3.0), Cm(7.6), Cm(5.0), "Lekcija", "Prompt nije validacija.\n\nI kad tekst izgleda hrvatski, može pasti na fonemima, duplikatima, HJP-u ili strukturi rečenice.", accent=BLUE, body_size=12.5)
    add_textbox(s, "Rješenje: model generira više kandidata nego što treba, a Python prihvaća samo one koji prođu pravila.", Cm(2.0), Cm(10.0), Cm(23.0), Cm(0.9), font_size=14, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    # 14 technical validity and PCD
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Kako evaluator odlučuje valjanost + PCD", 14)
    add_card(s, Cm(1.0), Cm(2.55), Cm(8.0), Cm(6.8), "Tehnički pragovi", "Kandidat je tehnički valjan ako prođe:\n• zasićenje ≥ traženi SL\n• samo hrvatska slova i razmaci\n• word = točno 1 riječ\n• sentence = 3–5 riječi\n• nije duplikat u runu\n• nema ponovljene riječi u rečenici", accent=TEAL, body_size=10.8)
    add_card(s, Cm(9.6), Cm(2.55), Cm(8.0), Cm(6.8), "Razina eksperimenta", "Korišteni pragovi:\n• paper riječi: SL 40, 50, 60, 70, 80\n• paper rečenice: SL 50 i 70\n• Task 16: SL 50 i 70\n\nSL = minimalni % fonema ciljne klase.", accent=ORANGE, body_size=10.8)
    add_card(s, Cm(18.2), Cm(2.55), Cm(8.0), Cm(6.8), "PCD", "Phonetic Content Dissimilarity mjeri koliko su dva kandidata fonetski različita.\n\nNaš paper-style PCD koristi dulji kandidat kao nazivnik.\nVeći PCD = veća raznolikost.", accent=BLUE, body_size=10.8)
    rows = [
        ["ChatGPT Plus", "0,474", "0,427", "0,360", "0,423", "0,487"],
        ["Ollama llama3.1:8b", "—", "0,537", "0,376", "0,522", "0,667"],
    ]
    add_table(
        s,
        Cm(2.0),
        Cm(10.1),
        Cm(23.2),
        Cm(2.0),
        ["Avg. PCD po klasama", "N", "SN", "S", "SV", "V"],
        rows,
        widths=[0.34, 0.132, 0.132, 0.132, 0.132, 0.132],
        font_size=9.2,
    )
    add_textbox(s, "— = nema dovoljno validnih kandidata za parni PCD u toj klasi.", Cm(2.2), Cm(12.2), Cm(22.8), Cm(0.5), font_size=8.8, color=MUTED, align=PP_ALIGN.CENTER)

    # 15 lexical review
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Leksička provjera: Hunspell screening + ručni HJP", 15)
    add_card(s, Cm(1.0), Cm(2.55), Cm(7.7), Cm(4.9), "Hunspell", "Automatski provjerava svaku normaliziranu riječ preko hrvatskog rječnika `hr_HR`.\n\nIzlaz: yes / no / unsure + popis nepoznatih ili odbačenih riječi.", accent=TEAL, body_size=11.2)
    add_card(s, Cm(9.7), Cm(2.55), Cm(7.7), Cm(4.9), "Što Hunspell nije", "Nije HJP dokaz.\nNe provjerava značenje rečenice.\nNe provjerava kliničku prikladnost.\nMože odbiti valjane oblike ili prihvatiti kontekstualno loše riječi.", accent=RED, body_size=11.2)
    add_card(s, Cm(18.4), Cm(2.55), Cm(7.7), Cm(4.9), "Ručni HJP review", "Izvoz jedinstvenih riječi → ručno `hjp_valid`.\n\nRiječ: validna ako je ta riječ yes.\nRečenica: validna samo ako su sve riječi yes.", accent=ORANGE, body_size=11.2)
    add_metric(s, Cm(1.4), Cm(8.25), Cm(5.3), Cm(2.0), "95,4 %", "ChatGPT Plus Hunspell-valid", color=BLUE)
    add_metric(s, Cm(7.5), Cm(8.25), Cm(5.3), Cm(2.0), "776 / 797", "HJP-valid nakon word-review", color=GREEN)
    add_metric(s, Cm(13.6), Cm(8.25), Cm(5.3), Cm(2.0), "21", "HJP-invalid kandidata", color=RED)
    add_metric(s, Cm(19.7), Cm(8.25), Cm(5.3), Cm(2.0), "355", "TTS-ready: tehnički + HJP", color=TEAL)

    # 16 what affects results
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Što točno utječe na rezultate?", 16)
    rows = [
        ["SL prolaznost", "samo fonemi", "target_count / total ≥ SL", "glavni fonetski kriterij"],
        ["Tehnička valjanost", "struktura + fonemi", "SL, znakovi, broj riječi, duplikat, ponavljanje", "osnovni Python pass/fail"],
        ["Hunspell valid", "leksika", "svaka riječ prihvaćena u hr_HR rječniku", "automatski screening"],
        ["Technical + Hunspell", "tehnika + screening", "tehnički validno i Hunspell=yes", "skalabilna leksička procjena"],
        ["HJP/manual valid", "ručni word review", "sve riječi u kandidatu imaju hjp_valid=yes", "najbliže kriteriju iz paper-a"],
        ["PCD", "raznolikost", "fonemska udaljenost među kandidatima", "nije pass/fail nego diversity"],
        ["TTS success", "audio tehnika", "WAV postoji, mono, 16 kHz, 16-bit PCM", "ne dokazuje prirodnost"],
        ["WER/CER + slušanje", "audio razumljivost", "ASR pogreške + ljudske ocjene", "proxy + subjektivna provjera"],
    ]
    add_table(
        s,
        Cm(0.9),
        Cm(2.55),
        Cm(25.4),
        Cm(7.8),
        ["Rezultat", "Što mjeri", "Što ga ruši / smanjuje", "Kako se koristi"],
        rows,
        widths=[0.18, 0.20, 0.36, 0.26],
        font_size=7.4,
    )
    add_card(
        s,
        Cm(1.5),
        Cm(11.0),
        Cm(24.2),
        Cm(1.5),
        "Najvažnija interpretacija",
        "U radu ne postoji jedna jedina “valjanost”. Rezultate čitam slojevito: prvo fonetski i strukturno, zatim automatski leksički, zatim ručno HJP, a tek onda audio.",
        accent=BLUE,
        body_size=11.5,
    )

    # 17 audio design
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Audio faza: isti tekst, tri TTS sustava", 17)
    add_card(s, Cm(1.2), Cm(3.0), Cm(7.4), Cm(4.6), "Ulaz", "Samo kandidati koji su prošli:\n• tehničku validaciju\n• Hunspell screening\n• ručni HJP word-review", accent=TEAL, body_size=12.5)
    add_card(s, Cm(9.9), Cm(3.0), Cm(7.4), Cm(4.6), "TTS sustavi", "• eSpeak NG, hrvatski glas hr\n• Coqui VITS HR\n• SpeechT5 HR lokalni model\n\nSvi dobivaju isti skup tekstova.", accent=ORANGE, body_size=12.5)
    add_card(s, Cm(18.6), Cm(3.0), Cm(7.4), Cm(4.6), "Izlaz", "Svi WAV-ovi normalizirani su na:\n• mono\n• 16 kHz\n• 16-bit PCM\n\nZatim ide ASR i slušna provjera.", accent=GREEN, body_size=12.5)
    add_textbox(s, "Uravnoteženi TTS podskup: 96 kandidata; skupina S / 70 % / riječ imala je samo 1 dostupni validirani kandidat.", Cm(1.8), Cm(9.7), Cm(23.6), Cm(0.9), font_size=12.5, color=MUTED, align=PP_ALIGN.CENTER)

    # 18 TTS and ASR
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Rezultat 3: tehnička TTS provjera i ASR WER/CER", 18)
    rows = [
        ["Coqui VITS HR", "96/96", "100 %", "WER 1,187", "CER 0,693"],
        ["eSpeak NG", "96/96", "100 %", "WER 1,076", "CER 0,443"],
        ["SpeechT5 HR", "96/96", "100 %", "WER 0,967", "CER 0,438"],
    ]
    add_table(s, Cm(1.2), Cm(3.0), Cm(24.8), Cm(3.4), ["TTS", "sintetizirano", "format OK", "ASR WER", "ASR CER"], rows, widths=[0.28, 0.18, 0.18, 0.18, 0.18], font_size=10)
    add_bar_chart(s, Cm(2.0), Cm(7.3), Cm(10.5), Cm(4.6), ["Coqui", "eSpeak", "SpeechT5"], [("WER ×100", [118.7, 107.6, 96.7]), ("CER ×100", [69.3, 44.3, 43.8])], title="Niže je bolje; skala radi prikaza ×100")
    add_card(s, Cm(14.0), Cm(7.4), Cm(11.2), Cm(3.8), "Caveat", "ASR WER/CER je relativni proxy za razumljivost, nije klinički dokaz izgovora. Visok WER može doći od TTS-a, ASR-a ili neobičnog fonemskog sastava materijala.", accent=RED, body_size=12)

    # 19 audio demo
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Audio demo: ista rečenica kroz tri TTS sustava", 19)
    add_card(s, Cm(1.4), Cm(2.8), Cm(24.6), Cm(2.0), "Tekst za usporedbu", "Sva tri audio primjera koriste isti kandidat: „Draga rada radi.”\nKopije su pojačane samo za prezentaciju, posebno zbog tišeg SpeechT5 izlaza.", accent=BLUE, body_size=13)
    x_positions = [Cm(1.8), Cm(10.0), Cm(18.2)]
    colors = [GREEN, ORANGE, TEAL]
    for idx, (label, path) in enumerate(audio_demo.items(), start=1):
        x = x_positions[idx - 1]
        add_metric(s, x, Cm(6.0), Cm(6.5), Cm(2.2), str(idx), label, color=colors[idx - 1])
        file_label = path.name if path.exists() else "audio datoteka nije pronađena"
        box = add_textbox(s, file_label, x, Cm(8.65), Cm(6.5), Cm(0.8), font_size=9.2, color=BLUE, align=PP_ALIGN.CENTER)
        if path.exists():
            box.click_action.hyperlink.address = f"audio_demo/{path.name}"
        add_textbox(s, "Klikni naziv datoteke ili pusti WAV iz mape audio_demo.", x, Cm(9.55), Cm(6.5), Cm(0.8), font_size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, f"Mapa s demo snimkama: {AUDIO_DEMO_DIR}", Cm(1.5), Cm(12.3), Cm(24.0), Cm(0.7), font_size=9.5, color=MUTED, align=PP_ALIGN.CENTER)

    # 20 human listening
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Rezultat 4: slušna provjera na uravnoteženom uzorku", 20)
    rows = [
        ["Coqui VITS HR", "20", "2,15", "2,10", "1,233"],
        ["eSpeak NG", "20", "4,65", "3,30", "0,967"],
        ["SpeechT5 HR", "20", "3,35", "3,60", "0,958"],
    ]
    add_table(s, Cm(1.2), Cm(3.0), Cm(24.8), Cm(3.4), ["TTS", "N", "razumljivost /5", "prirodnost /5", "ASR WER uzorka"], rows, widths=[0.30, 0.12, 0.21, 0.21, 0.16], font_size=10)
    add_card(s, Cm(1.5), Cm(7.5), Cm(7.5), Cm(3.8), "Najrazumljiviji", "eSpeak NG\nprosjek 4,65/5\nmedijan 5/5", accent=GREEN, body_size=14)
    add_card(s, Cm(9.9), Cm(7.5), Cm(7.5), Cm(3.8), "Najprirodniji", "SpeechT5 HR\nprosjek 3,60/5\nmedijan 4/5", accent=TEAL, body_size=14)
    add_card(s, Cm(18.3), Cm(7.5), Cm(7.5), Cm(3.8), "Najslabiji", "Coqui VITS HR\nniža razumljivost i prirodnost\nu ovom materijalu", accent=RED, body_size=14)

    # 21 achieved/proxy/future
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, "Što je ostvareno, što nije i sljedeći koraci", 21)
    add_card(s, Cm(1.0), Cm(2.7), Cm(7.7), Cm(6.2), "Ostvareno", "• reprodukcija ideje iz literature\n• deterministički hrvatski fonemizator\n• CSV/Markdown izvještaji i PCD\n• ChatGPT vs. Ollama usporedba\n• HJP word-review\n• 3 TTS sustava + ASR + slušna provjera", accent=GREEN, body_size=11.5)
    add_card(s, Cm(9.7), Cm(2.7), Cm(7.7), Cm(6.2), "Nije ostvareno / ograničenja", "• nije klinički odobren materijal\n• ChatGPT Plus je ručni, ne API\n• HJP nije automatski scrap-an\n• slušna provjera ima jednog recenzenta\n• ASR je proxy, ne dokaz izgovora", accent=ORANGE, body_size=11.5)
    add_card(s, Cm(18.4), Cm(2.7), Cm(7.7), Cm(6.2), "Problemi i rješenja", "• LLM ne zna pouzdano brojati foneme → Python validacija\n• duplikati → detekcija i izvještaj\n• lokalni modeli slabi → odvojena usporedba\n• TTS formati različiti → normalizacija WAV-a", accent=BLUE, body_size=11.5)
    add_textbox(s, "Zaključak: pipeline je reproducibilan i pokazuje da generiranje materijala može pomoći, ali konačna lingvistička i klinička odluka ostaje na čovjeku.", Cm(1.8), Cm(10.2), Cm(23.6), Cm(1.1), font_size=14, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    prs.save(PPTX_OUT)
    NOTES_OUT.write_text(build_notes(), encoding="utf-8")


def build_notes() -> str:
    return """# Govorne bilješke za prezentaciju

Trajanje: ciljaj 12-14 minuta i ostavi 1-2 minute za pitanja. Govori mirno; gotovo svaka slajd ima jednu glavnu poruku.

## Slajd 1 — Naslov

Reci: Tema je evaluacija jezičnih generativnih modela za audiorehabilitaciju. Cilj nije samo generirati lijepe riječi, nego provjeriti mogu li se dobiti hrvatske riječi i kratke rečenice s kontroliranim fonemskim sastavom, a zatim ih pretvoriti u audio.

## Slajd 2 — Motivacija

Objasni problem: u audiorehabilitaciji trebamo materijal koji cilja određene glasove ili kontraste. Ako korisnik ima problem razlikovati /s/ i /š/, nije dovoljno imati bilo koju rečenicu. Treba znati koji fonemi se stvarno pojavljuju i u kojem omjeru.

## Slajd 3 — Istraživačka pitanja i doprinosi

Glavna ideja: rad nije samo LLM generiranje. Ima četiri sloja: LLM predlaže hrvatske kandidate, Python ih deterministički validira, TTS pretvara samo validirane kandidate u audio, a ASR i slušna provjera daju dodatnu audio evaluaciju.

Doprinos rada je reproducibilni pipeline: LLM/CSV kandidati → hrvatski fonemizator → deterministička validacija → Hunspell/HJP review → TTS usporedba → ASR WER/CER + slušna provjera.

Ako profesor pita „zašto ne vjerovati ChatGPT-u?”: Zato što LLM često daje uvjerljiv odgovor, ali nema garantirano točno brojanje fonema, pogotovo za hrvatske višeslovne foneme dž, lj i nj.

## Slajd 4 — Veza s literaturom

Reci da se oslanjamo na rad Andrijašević i Vukelić, gdje se GPT koristi za generiranje hrvatskog govornog materijala za auditivni trening. Mi preuzimamo ideju pet klasa fonema i razina zasićenja, ali uvodimo strožu validaciju.

Važno: HJP u promptu je samo uputa modelu. U našem radu HJP valjanost je zasebna ručna provjera riječi.

## Slajd 5 — Fonemske klase i formula

Objasni formulu: broj fonema ciljne klase dijeli se s ukupnim brojem fonema i množi sa 100. Ako ciljamo klasu N na 70 %, kandidat prolazi samo ako barem 70 % njegovih fonema pripada klasi N.

Tehnički detalj: parser ide longest-match-first za dž, lj i nj. Bez toga bi npr. „panj” pogrešno postao p-a-n-j umjesto p-a-nj.

## Slajd 6 — Pipeline

Objasni što skripta radi:

1. Učita CSV ili generira kandidate preko Ollame.
2. Normalizira tekst u UTF-8, čuva hrvatske znakove.
3. Fonemizira tekst.
4. Računa zasićenje.
5. Provjerava dopuštene znakove, broj riječi, duplikate i ponovljene riječi.
6. Dodaje Hunspell screening i ručni HJP word-review.
7. Stvara CSV i Markdown izvještaje.

Ako pita „koje su glavne datoteke”: `src/phonemizer.py`, `src/validators.py`, `src/metrics.py`, `src/pipeline.py`, `src/tts.py`, `src/asr_eval.py`.

## Slajd 7 — Tehnička implementacija

Ovdje pokaži da je rad računalni sustav, ne samo promptanje. Ukratko prođi module.

`phonemizer.py` normalizira tekst i pretvara ga u foneme. `phoneme_classes.py` zna kojoj klasi pripada svaki fonem. `validators.py` računa saturation i dodaje failure reasons. `metrics.py` računa PCD, frekvencije i duplicate rate. `generators.py` služi za ručni CSV i Ollama generiranje. `pipeline.py` sve povezuje preko CLI naredbi. `tts.py` i `asr_eval.py` rade audio fazu.

Naglasak: svi rezultati su CSV/Markdown, UTF-8, s `run_id`, tako da je eksperiment ponovljiv.

## Slajd 8 — Algoritam

Objasni algoritam jednostavno: prvo normalizacija teksta, zatim parsiranje fonema. Najvažnije je longest-match-first za `dž`, `lj`, `nj`, jer su to jedan fonem, a ne dva slova.

Nakon toga se izračuna saturation: broj fonema ciljne klase podijeljen s ukupnim brojem fonema. Ako je prag 70 %, kandidat s 65 % pada, bez obzira na to koliko lijepo izgleda.

Failure reasons su bitni jer kandidat može pasti iz više razloga. Na primjer, rečenica može istovremeno imati premalo zasićenje i pogrešan broj riječi.

## Slajd 9 — Dizajn eksperimenata

Reci da postoje tri sloja:

Prvi je reprodukcija paper-style promptova. Drugi je Task 16 usporedba ChatGPT Plus protiv lokalne Ollame. Treći je audio faza samo na validiranim kandidatima.

Naglasak: audio ne radimo nad svim kandidatima, nego samo nad onima koji su prošli determinističku i leksičku provjeru.

## Slajd 10 — Reprodukcija rada: riječi

Ovdje reci da je ovo namjerno prikazano gotovo istim formatom kao Table II u radu. U paperu su retci fonemske klase, stupci su saturation leveli, a ćelija je postotak riječi koje zadovoljavaju oba kriterija: zasićenje + standardni hrvatski/HJP.

Naša tablica koristi isti princip: saturation pass + ručni HJP word-review. Važno je reći da naš pipeline dodatno mjeri duplikate, ali ih ovdje ne miješamo u ćeliju jer želimo metodološki isti indikator kao u paperu.

Zaključak za riječi: naši paper-style rezultati su vrlo visoki za većinu klasa, ali SV na 70 % i V na 80 % pokazuju pad. To je usporedivo s idejom rada da visoka zasićenja i “visoke” klase postaju teže.

## Slajd 11 — Reprodukcija rada: rečenice

Ovdje je analog Table III iz rada. Paper za rečenice prikazuje postotak rečenica koje zadovoljavaju saturation level. Mi prikazujemo isti kriterij.

Glavna poruka: na 50 % smo gotovo savršeni za većinu klasa. Na 70 % se vidi pad, posebno za klasu V. To je važan zaključak jer se slaže s očekivanjem: viši prag daje modelu manje slobode.

## Slajd 12 — ChatGPT Plus vs Ollama

Ovo je najvažniji tekstni rezultat. ChatGPT Plus ima 48,6 % tehnički valjanih kandidata i 85,9 % saturation pass. Ollama ima samo 4,5 % tehnički valjanih i 4,6 % saturation pass.

Zaključak: u ovoj konfiguraciji lokalni `llama3.1:8b` nije dovoljan za strogo hrvatsko fonemsko zasićenje. ChatGPT Plus je znatno bolji generator, ali i dalje treba validaciju.

## Slajd 13 — Analiza grešaka

Objasni razliku u vrstama grešaka:

ChatGPT Plus: najveći problem su duplikati. To znači da model često razumije zadatak, ali se vrti oko istih sigurnih kandidata.

Ollama: najveći problem je sam fonetski kriterij. To je ozbiljniji problem jer znači da generirani tekst ne zadovoljava glavnu znanstvenu mjeru.

## Slajd 14 — Tehnička valjanost i PCD

Ovo je slajd za metodološko objašnjenje. Reci da `is_valid` nije procjena LLM-a, nego rezultat determinističkih pravila.

Kandidat je tehnički valjan ako prođe sve ove pragove: zasićenje mora biti veće ili jednako traženom SL-u, znakovi smiju biti samo hrvatska slova i razmaci, riječ mora imati točno jednu riječ, rečenica mora imati 3 do 5 riječi, kandidat ne smije biti duplikat u istom runu i rečenica ne smije ponavljati istu riječ.

U eksperimentima su pragovi bili: za paper riječi 40, 50, 60, 70 i 80 %, za paper rečenice 50 i 70 %, a za Task 16 usporedbu 50 i 70 %. Dakle `saturation_level` nije procjena, nego konkretan numerički prag u CSV-u.

PCD znači Phonetic Content Dissimilarity. Naš paper-style PCD uspoređuje fonemski sadržaj dvaju kandidata: gleda koliko fonema u duljem kandidatu nije dijeljeno s kraćim, podijeljeno duljinom duljeg kandidata. Veći prosječni PCD znači raznolikiji skup.

Rezultat: ChatGPT Plus ima viši prosječni PCD po grupama od Ollame, što znači da je osim više validnih kandidata dao i korisniju fonetsku raznolikost. Ollama ima malo grupa s dovoljno validnih kandidata, pa je i PCD slabiji.

## Slajd 15 — Hunspell i HJP

Reci da Hunspell nije HJP. Hunspell je automatski spelling/dictionary screening: pipeline uzima normalizirane riječi i šalje ih lokalnom alatu `hunspell -d hr_HR -l`. Ako Hunspell vrati riječ kao nepoznatu, kandidat dobiva `dictionary_failed`, a u CSV-u se sprema `dictionary_invalid_words`.

Važno: Hunspell radi riječ-po-riječ. On ne zna je li rečenica semantički prirodna, klinički prikladna ili stvarno potvrđena na HJP-u. Može odbaciti valjani flektirani oblik, a može i prihvatiti riječ koja nije dobra za rehabilitacijski materijal.

Ručni HJP word-review: iz kandidata se izvezu jedinstvene riječi, označi se `hjp_valid`, a zatim se ta odluka propagira natrag na riječi i rečenice. Rečenica je HJP-valid samo ako su sve njezine riječi HJP-valid.

## Slajd 16 — Što utječe na rezultate

Ovaj slajd služi da ne pomiješamo različite razine validacije.

Prvo je `SL prolaznost`: to je samo fonetski kriterij. Ako kandidat nema dovoljno fonema ciljne klase, pada bez obzira na to je li riječ hrvatska.

Drugo je `tehnička valjanost`: u normalnoj tekstnoj usporedbi to znači saturation pass, dopušteni znakovi, ispravan broj riječi, bez duplikata i bez ponavljanja riječi. To je osnovni Python pass/fail.

Treće je Hunspell: to je automatski leksički screening. On utječe na Hunspell valid rate i technical + Hunspell valid rate. Ne smije se interpretirati kao HJP dokaz.

Četvrto je HJP/manual valid: to je ručni word-level review. Za riječ je dovoljno da je ta riječ yes. Za rečenicu sve riječi moraju biti yes. Ovo je najbliže kriteriju iz paper-a o standardnom hrvatskom jeziku.

PCD ne odlučuje valjanost. On samo govori koliko su kandidati fonetski raznoliki. TTS success govori samo je li audio tehnički napravljen. WER/CER i slušanje govore o audio razumljivosti, ali nisu kliničko odobrenje.

## Slajd 17 — Audio dizajn

Objasni da je za poštenu TTS usporedbu isti skup od 96 tekstova poslan u sva tri sustava: eSpeak NG, Coqui VITS HR i SpeechT5 HR. Time se ne uspoređuju različiti tekstovi, nego različiti sintetizatori na istom materijalu.

Svi izlazi su normalizirani u isti format: WAV, mono, 16 kHz, 16-bit PCM.

## Slajd 18 — TTS i ASR

Sva tri TTS sustava tehnički su uspjela sintetizirati 96/96 kandidata. ASR evaluacija koristi fiksni `faster-whisper large-v3-turbo` profil za hrvatski.

Reci caveat: WER/CER nije klinički dokaz. To je automatski proxy. Visok WER može značiti da je TTS loš, ASR loš, ili da je materijal fonemski neobičan.

## Slajd 19 — Audio demo

Ovdje pusti tri WAV datoteke istim redom: eSpeak NG, Coqui VITS HR, SpeechT5 HR. Naglasi da je tekst isti: „Draga rada radi.” Time se sluša razlika u sintetizatoru, a ne razlika u sadržaju.

Reci da su demo kopije pojačane za prezentaciju jer je SpeechT5 u originalnom izlazu bio tiši. Originalni istraživački audio nije mijenjan; pojačanje je samo praktično za slušanje u učionici.

## Slajd 20 — Slušna provjera

Ovo je važan ljudski rezultat. eSpeak NG je najrazumljiviji po ljudskoj ocjeni, iako nije najprirodniji. SpeechT5 je najprirodniji, ali ne najrazumljiviji. Coqui je u ovom testu najslabiji.

Poanta: automatska ASR metrika i ljudska procjena ne moraju dati isti poredak, zato se ne smije stati samo na WER/CER.

## Slajd 21 — Zaključak

Zaključi:

1. Pipeline radi i reproducibilan je.
2. LLM može pomoći u generiranju kandidata, ali samo ako Python deterministički validira.
3. ChatGPT Plus je bio bolji od lokalne Ollame za tekst.
4. Audio faza je tehnički izvediva, ali klinička prikladnost traži stručnu provjeru.

## Kratki odgovori ako profesor pita

**Kako fonemizator radi?**  
Tekst se pretvara u mala slova, uklanja se interpunkcija, čuvaju se hrvatski znakovi, a `dž`, `lj` i `nj` se prepoznaju prije pojedinačnih slova. Razmaci ne ulaze u broj fonema.

**Što znači saturation pass?**  
Ako kandidat ima N fonema, a nk pripada ciljnoj klasi, zasićenje je nk/N × 100. Kandidat prolazi ako je postotak veći ili jednak traženom pragu.

**Zašto imate Hunspell i HJP?**  
Hunspell je automatski i skalabilan, ali nije konačan. HJP/manual word-review je bliže kriteriju iz literature, ali je ručni.

**Zašto je Ollama tako slaba?**  
Testirani lokalni model nije pouzdano pratio hrvatske fonemske klase. Većina kandidata pala je na zasićenju, ne samo na hrvatskoj leksici.

**Zašto TTS nije klinički zaključak?**  
Jer tehnički WAV i dobra ASR transkripcija ne dokazuju da je izgovor terapeutski prikladan. Potrebni su stručnjaci, više slušatelja i eventualno korisničko testiranje.

**Što biste poboljšali?**  
Veći ručni listening review, više ocjenjivača, inter-rater reliability, jači lokalni LLM, repair loop za neuspjele kandidate i formalniji HJP/klinički review.
"""


if __name__ == "__main__":
    build_deck()
    print(PPTX_OUT)
    print(NOTES_OUT)
